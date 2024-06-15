import streamlit as st
from pptx import Presentation
from docx import Document
from io import BytesIO
from docx.shared import Pt
import re

def clean_text(text):
    # Remove control characters
    cleaned_text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
    return cleaned_text

def ppt_to_word(ppt_file):
    # Load the presentation
    presentation = Presentation(ppt_file)

    # Create a Word document
    doc = Document()

    # Extract text from slides and add to the Word document
    slide_texts = []
    for slide in presentation.slides:
        slide_text = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                cleaned_paragraph_text = clean_text(paragraph.text)
                slide_text += cleaned_paragraph_text + "\n"
                p = doc.add_paragraph(cleaned_paragraph_text)
                p.style.font.size = Pt(12)
        slide_texts.append(slide_text)

    word_file = BytesIO()
    doc.save(word_file)
    word_file.seek(0)
    return word_file, slide_texts

def generate_questions(slide_texts):
    questions = []
    for index, slide_text in enumerate(slide_texts):
        lines = slide_text.strip().split('\n')
        for i, line in enumerate(lines):
            if line:  # Skip empty lines
                question = f"Slide {index+1}, Point {i+1}: What does this mean?\n\n{line}"
                questions.append(question)
    return questions

def main():
    st.title("PPT to Word Converter")

    uploaded_file = st.file_uploader("Choose a PPT file", type="pptx")

    if uploaded_file is not None:
        if st.button("Convert"):
            with st.spinner("Converting..."):
                word_file, slide_texts = ppt_to_word(uploaded_file)

            st.success("Conversion successful!")

            st.download_button(
                label="Download Word file",
                data=word_file,
                file_name="converted.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )

            questions = generate_questions(slide_texts)
            st.header("Generated Questions")
            for question in questions:
                st.write(question)
                user_response = st.text_area(f"Your response to: {question}", key=question)
                if user_response:
                    st.write(f"Response recorded: {user_response}")

if __name__ == "__main__":
    main()
